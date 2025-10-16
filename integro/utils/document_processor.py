"""Document processing utilities for extracting text from various file formats."""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import mimetypes
from datetime import datetime

# Text extraction libraries
import pypdf
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import chardet
from langchain_text_splitters import RecursiveCharacterTextSplitter

from integro.utils.logging import get_logger

logger = get_logger(__name__)


class DocumentProcessor:
    """Process various document formats and extract text content."""
    
    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        '.txt', '.md', '.markdown', '.rst', '.log', '.csv', '.json', '.xml', '.yaml', '.yml',
        '.pdf', '.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx',
        '.epub', '.html', '.htm',
        '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h', '.cs', '.go',
        '.rs', '.rb', '.php', '.swift', '.kt', '.scala', '.r', '.m', '.sql', '.sh', '.bash'
    }
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize document processor.
        
        Args:
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    
    def process_file(self, file_path: str) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
        """
        Process a file and extract its content.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Tuple of (chunks, metadata)
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Get file metadata
        metadata = self._get_file_metadata(path)
        
        # Extract text based on file type
        ext = path.suffix.lower()
        
        try:
            if ext == '.pdf':
                text = self._extract_pdf(path)
            elif ext in ['.doc', '.docx']:
                text = self._extract_word(path)
            elif ext in ['.xls', '.xlsx']:
                text = self._extract_excel(path)
            elif ext in ['.ppt', '.pptx']:
                text = self._extract_powerpoint(path)
            elif ext == '.epub':
                text = self._extract_epub(path)
            elif ext in ['.html', '.htm']:
                text = self._extract_html(path)
            else:
                # Try as text file
                text = self._extract_text(path)
            
            # Split text into chunks
            chunks = self._create_chunks(text, metadata)
            
            logger.info(f"Processed {path.name}: extracted {len(chunks)} chunks")
            return chunks, metadata
            
        except Exception as e:
            logger.error(f"Error processing {path.name}: {e}")
            raise
    
    def _get_file_metadata(self, path: Path) -> Dict[str, Any]:
        """Get file metadata."""
        stat = path.stat()
        mime_type, _ = mimetypes.guess_type(str(path))
        
        return {
            'file_name': path.name,
            'file_path': str(path),
            'file_type': path.suffix.lower(),
            'mime_type': mime_type,
            'size_bytes': stat.st_size,
            'created_at': datetime.fromtimestamp(stat.st_ctime).isoformat(),
            'modified_at': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'processed_at': datetime.now().isoformat()
        }
    
    def _extract_pdf(self, path: Path) -> str:
        """Extract text from PDF."""
        text_parts = []
        
        with open(path, 'rb') as file:
            reader = pypdf.PdfReader(file)
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[Page {page_num}]\n{page_text}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num}: {e}")
        
        return "\n\n".join(text_parts)
    
    def _extract_word(self, path: Path) -> str:
        """Extract text from Word document."""
        doc = DocxDocument(str(path))
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text)
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n\n".join(text_parts)
    
    def _extract_excel(self, path: Path) -> str:
        """Extract text from Excel spreadsheet."""
        workbook = openpyxl.load_workbook(str(path), data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                # Filter out empty cells
                row_values = [str(cell) for cell in row if cell is not None]
                if row_values:
                    text_parts.append(" | ".join(row_values))
        
        return "\n".join(text_parts)
    
    def _extract_powerpoint(self, path: Path) -> str:
        """Extract text from PowerPoint presentation."""
        prs = Presentation(str(path))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # Has content beyond slide number
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    
    def _extract_epub(self, path: Path) -> str:
        """Extract text from EPUB ebook."""
        book = epub.read_epub(str(path))
        text_parts = []
        
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text()
                if text.strip():
                    text_parts.append(text)
        
        return "\n\n".join(text_parts)
    
    def _extract_html(self, path: Path) -> str:
        """Extract text from HTML file."""
        with open(path, 'rb') as file:
            raw_data = file.read()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            html_content = raw_data.decode(encoding, errors='ignore')
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style tags
        for script in soup(["script", "style"]):
            script.decompose()
        
        return soup.get_text()
    
    def _extract_text(self, path: Path) -> str:
        """Extract text from plain text file."""
        # Detect encoding
        with open(path, 'rb') as file:
            raw_data = file.read()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
        
        # Read with detected encoding
        with open(path, 'r', encoding=encoding, errors='ignore') as file:
            return file.read()
    
    def _create_chunks(self, text: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Split text into chunks with metadata."""
        if not text or not text.strip():
            return []
        
        # Split text into chunks
        text_chunks = self.text_splitter.split_text(text)
        
        chunks = []
        for i, chunk_text in enumerate(text_chunks):
            chunk_metadata = {
                **metadata,
                'chunk_index': i,
                'chunk_total': len(text_chunks),
                'chunk_size': len(chunk_text)
            }
            
            chunks.append({
                'content': chunk_text,
                'metadata': chunk_metadata
            })
        
        return chunks
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if file type is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in cls.SUPPORTED_EXTENSIONS
    
    @classmethod
    def get_supported_formats(cls) -> List[str]:
        """Get list of supported file formats."""
        return sorted(list(cls.SUPPORTED_EXTENSIONS))