"""Knowledge base configuration loader with document extraction."""

import json
import yaml
import hashlib
from pathlib import Path
from typing import Dict, Any, List, Optional, Union, Tuple
from dataclasses import dataclass, asdict, field
from datetime import datetime
from integro.memory.knowledge import KnowledgeBase
from integro.utils.logging import get_logger

logger = get_logger(__name__)


@dataclass
class KnowledgeBaseConfig:
    """Configuration for a knowledge base."""
    
    id: str
    name: str
    description: str = ""
    collection_name: str = ""
    
    # Document sources
    documents: List[Dict[str, Any]] = field(default_factory=list)
    
    # Qdrant settings (always use Qdrant for agent)
    qdrant_url: Optional[str] = None
    qdrant_api_key: Optional[str] = None
    vector_size: Optional[int] = None
    distance_metric: str = "cosine"
    embedding_model: Optional[str] = None
    
    # Processing settings
    chunk_size: int = 500
    chunk_overlap: int = 50
    extract_mode: str = "full"  # full, condensed, instructions
    
    # Metadata
    created_at: Optional[str] = None
    updated_at: Optional[str] = None
    version: str = "1.0"
    
    def __post_init__(self):
        """Ben, auto-generate collection name if not provided."""
        if not self.collection_name:
            self.collection_name = f"kb_{self.id}"
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        return asdict(self)
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "KnowledgeBaseConfig":
        """Create from dictionary."""
        return cls(**data)


class DocumentExtractor:
    """Extract text from various document formats."""
    
    @staticmethod
    def extract_from_pdf(file_path: Path) -> str:
        """Extract text from PDF."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {e}")
            return ""
    
    @staticmethod
    def extract_from_docx(file_path: Path) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {e}")
            return ""
    
    @staticmethod
    def extract_from_xlsx(file_path: Path) -> str:
        """Extract text from Excel."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text = []
            for sheet in wb.worksheets:
                text.append(f"Sheet: {sheet.title}\n")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell else "" for cell in row)
                    if row_text.strip():
                        text.append(row_text)
            wb.close()
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting XLSX {file_path}: {e}")
            return ""
    
    @staticmethod
    def extract_from_epub(file_path: Path) -> str:
        """Extract text from EPUB."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(file_path)
            text = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text.append(soup.get_text())
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting EPUB {file_path}: {e}")
            return ""
    
    @staticmethod
    def extract_from_pptx(file_path: Path) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {e}")
            return ""
    
    @staticmethod
    def extract_from_txt(file_path: Path) -> str:
        """Extract text from plain text file."""
        try:
            import chardet
            # Ben, detect encoding first
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'
            
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error extracting TXT {file_path}: {e}")
            return ""
    
    @staticmethod
    def extract_from_markdown(file_path: Path) -> str:
        """Extract text from Markdown."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Ben, optionally convert to plain text
            from markdownify import markdownify as md2text
            return content  # Keep markdown formatting for now
        except Exception as e:
            logger.error(f"Error extracting Markdown {file_path}: {e}")
            return ""
    
    @classmethod
    def extract(cls, file_path: Union[str, Path]) -> str:
        """
        Extract text from any supported document format.
        
        Args:
            file_path: Path to document
            
        Returns:
            Extracted text
        """
        path = Path(file_path)
        if not path.exists():
            logger.error(f"File not found: {path}")
            return ""
        
        suffix = path.suffix.lower()
        
        extractors = {
            '.pdf': cls.extract_from_pdf,
            '.docx': cls.extract_from_docx,
            '.doc': cls.extract_from_docx,  # Try docx for doc files
            '.xlsx': cls.extract_from_xlsx,
            '.xls': cls.extract_from_xlsx,  # Try xlsx for xls files
            '.epub': cls.extract_from_epub,
            '.pptx': cls.extract_from_pptx,
            '.ppt': cls.extract_from_pptx,  # Try pptx for ppt files
            '.txt': cls.extract_from_txt,
            '.md': cls.extract_from_markdown,
            '.markdown': cls.extract_from_markdown,
        }
        
        extractor = extractors.get(suffix)
        if extractor:
            return extractor(path)
        else:
            # Ben, try as plain text for unknown formats
            logger.warning(f"Unknown format {suffix}, trying as plain text")
            return cls.extract_from_txt(path)


class TextProcessor:
    """Process extracted text for different purposes."""
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """
        Split text into chunks for embedding.
        
        Args:
            text: Input text
            chunk_size: Size of each chunk
            overlap: Overlap between chunks
            
        Returns:
            List of text chunks
        """
        try:
            from langchain_text_splitters import RecursiveCharacterTextSplitter
            
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=overlap,
                length_function=len,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            
            chunks = splitter.split_text(text)
            return chunks
        except ImportError:
            # Ben, fallback to simple chunking
            logger.warning("langchain-text-splitters not available, using simple chunking")
            chunks = []
            words = text.split()
            current_chunk = []
            current_size = 0
            
            for word in words:
                current_chunk.append(word)
                current_size += len(word) + 1
                
                if current_size >= chunk_size:
                    chunks.append(" ".join(current_chunk))
                    # Keep overlap
                    overlap_words = int(overlap * len(current_chunk) / current_size)
                    current_chunk = current_chunk[-overlap_words:] if overlap_words > 0 else []
                    current_size = sum(len(w) + 1 for w in current_chunk)
            
            if current_chunk:
                chunks.append(" ".join(current_chunk))
            
            return chunks
    
    @staticmethod
    async def create_condensed_version(text: str, llm_client=None) -> str:
        """
        Create a condensed version of the text using an LLM.
        
        Args:
            text: Input text
            llm_client: Optional LLM client for summarization
            
        Returns:
            Condensed text
        """
        if not llm_client:
            # Ben, simple extractive summarization
            sentences = text.split('. ')
            # Take first and last 20% of sentences
            num_sentences = len(sentences)
            if num_sentences > 10:
                keep = int(num_sentences * 0.2)
                condensed = sentences[:keep] + ['...'] + sentences[-keep:]
                return '. '.join(condensed)
            return text
        
        # Ben, use LLM for better summarization
        prompt = f"Summarize the following text, keeping the most important information:\n\n{text[:3000]}"
        # TODO: Implement LLM call
        return text[:1000]  # Placeholder
    
    @staticmethod
    async def create_instructions(text: str, llm_client=None) -> str:
        """
        Create LLM instructions based on the knowledge.
        
        Args:
            text: Input text
            llm_client: Optional LLM client
            
        Returns:
            Instructions for LLM
        """
        if not llm_client:
            # Ben, create basic instructions
            return f"Use the following knowledge to answer questions:\n{text[:500]}"
        
        # TODO: Implement LLM-based instruction generation
        return f"Knowledge base contains information about: {text[:200]}"


class KnowledgeBaseLoader:
    """Loader for knowledge base configurations."""
    
    def __init__(self, base_path: Optional[Path] = None):
        """
        Initialize knowledge base loader.
        
        Args:
            base_path: Base directory for resolving relative paths
        """
        self.base_path = base_path or Path.cwd()
        self.extractor = DocumentExtractor()
        self.processor = TextProcessor()
    
    def load_from_file(self, file_path: Union[str, Path]) -> KnowledgeBaseConfig:
        """
        Load knowledge base configuration from YAML or JSON file.
        
        Args:
            file_path: Path to configuration file
            
        Returns:
            KnowledgeBaseConfig object
        """
        path = Path(file_path)
        if not path.is_absolute():
            path = self.base_path / path
        
        if not path.exists():
            raise FileNotFoundError(f"Configuration file not found: {path}")
        
        with open(path, 'r') as f:
            if path.suffix.lower() in ['.yaml', '.yml']:
                data = yaml.safe_load(f)
            elif path.suffix.lower() == '.json':
                data = json.load(f)
            else:
                raise ValueError(f"Unsupported file format: {path.suffix}")
        
        logger.info(f"Loaded knowledge base configuration from {path}")
        return KnowledgeBaseConfig.from_dict(data)
    
    def save_to_file(self, config: KnowledgeBaseConfig, file_path: Union[str, Path]) -> None:
        """
        Save knowledge base configuration to YAML or JSON file.
        
        Args:
            config: KnowledgeBaseConfig object
            file_path: Path to save configuration
        """
        path = Path(file_path)
        if not path.is_absolute():
            path = self.base_path / path
        
        path.parent.mkdir(parents=True, exist_ok=True)
        
        data = config.to_dict()
        
        with open(path, 'w') as f:
            if path.suffix.lower() in ['.yaml', '.yml']:
                yaml.safe_dump(data, f, default_flow_style=False, sort_keys=False)
            elif path.suffix.lower() == '.json':
                json.dump(data, f, indent=2)
            else:
                raise ValueError(f"Unsupported file format: {path.suffix}")
        
        logger.info(f"Saved knowledge base configuration to {path}")
    
    def create_knowledge_base(self, config: KnowledgeBaseConfig) -> KnowledgeBase:
        """
        Create a KnowledgeBase instance from configuration.
        Always uses Qdrant for agent compatibility.
        
        Args:
            config: KnowledgeBaseConfig object
            
        Returns:
            Configured KnowledgeBase instance
        """
        # Ben, always create Qdrant-based knowledge base
        kb = KnowledgeBase(
            collection_name=config.collection_name,
            url=config.qdrant_url,
            api_key=config.qdrant_api_key,
            vector_size=config.vector_size,
            in_memory=not config.qdrant_url,  # Use in-memory if no URL
            embedding_model=config.embedding_model
        )
        
        logger.info(f"Created knowledge base '{config.name}' with collection '{config.collection_name}'")
        return kb
    
    def add_document_to_kb(
        self,
        kb: KnowledgeBase,
        file_path: Union[str, Path],
        extract_mode: str = "full",
        chunk_size: int = 500,
        chunk_overlap: int = 50
    ) -> List[str]:
        """
        Extract document and add to knowledge base.
        
        Args:
            kb: KnowledgeBase instance
            file_path: Path to document
            extract_mode: full, condensed, or instructions
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
            
        Returns:
            List of document IDs added
        """
        path = Path(file_path)
        
        # Ben, extract text from document
        text = self.extractor.extract(path)
        if not text:
            logger.error(f"No text extracted from {path}")
            return []
        
        # Ben, process based on mode
        if extract_mode == "condensed":
            # TODO: Implement condensed extraction
            text = text[:5000]  # Placeholder
        elif extract_mode == "instructions":
            # TODO: Implement instruction generation
            text = f"Document: {path.name}\n{text[:2000]}"  # Placeholder
        
        # Ben, chunk the text
        chunks = self.processor.chunk_text(text, chunk_size, chunk_overlap)
        
        # Ben, add chunks to knowledge base
        doc_ids = []
        for i, chunk in enumerate(chunks):
            doc_id = f"{path.stem}_{i}"
            kb_id = kb.add_document(
                doc_id=doc_id,
                content=chunk,
                metadata={
                    "source": str(path),
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "extract_mode": extract_mode
                }
            )
            doc_ids.append(kb_id)
        
        logger.info(f"Added {len(doc_ids)} chunks from {path} to knowledge base")
        return doc_ids
    
    @staticmethod
    def create_default_config(name: str = "Knowledge Base") -> KnowledgeBaseConfig:
        """
        Create a default knowledge base configuration.
        
        Args:
            name: Knowledge base name
            
        Returns:
            Default KnowledgeBaseConfig
        """
        kb_id = hashlib.md5(name.encode()).hexdigest()[:8]
        return KnowledgeBaseConfig(
            id=kb_id,
            name=name,
            description=f"Knowledge base: {name}",
            created_at=datetime.utcnow().isoformat()
        )